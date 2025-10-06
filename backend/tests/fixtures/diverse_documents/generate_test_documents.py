#!/usr/bin/env python3
"""
Generate test documents for diverse document testing.

This script creates various test documents with different characteristics
for comprehensive testing of the document processing system.
"""

import os
import json
from pathlib import Path
from typing import Dict, Any
import random
import string

# Try to import document creation libraries
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer, PageBreak
    REPORTLAB_AVAILABLE = True
except ImportError:
    print("Warning: reportlab not installed. PDF generation will be skipped.")
    REPORTLAB_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    print("Warning: python-docx not installed. DOCX generation will be skipped.")
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    print("Warning: python-pptx not installed. PPTX generation will be skipped.")
    PYTHON_PPTX_AVAILABLE = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    OPENPYXL_AVAILABLE = True
except ImportError:
    print("Warning: openpyxl not installed. XLSX generation will be skipped.")
    OPENPYXL_AVAILABLE = False


BASE_DIR = Path(__file__).parent


def save_metadata(filepath: str, metadata: Dict[str, Any]):
    """Save document metadata to .info.json file."""
    info_file = f"{filepath}.info.json"
    with open(info_file, 'w') as f:
        json.dump(metadata, f, indent=2)


def generate_clean_pdf():
    """Generate clean digital PDF with tables and headings."""
    if not REPORTLAB_AVAILABLE:
        print("Skipping PDF generation - reportlab not available")
        return

    output_dir = BASE_DIR / "pdf" / "clean_digital"
    output_file = output_dir / "business_report_01.pdf"

    doc = SimpleDocTemplate(str(output_file), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Title
    story.append(Paragraph("Quarterly Business Report", styles['Title']))
    story.append(Spacer(1, 12))

    # Introduction
    story.append(Paragraph("Executive Summary", styles['Heading1']))
    story.append(Paragraph(
        "This report provides an overview of the company's performance in Q3 2025. "
        "Key metrics show strong growth across all departments with particular success "
        "in digital transformation initiatives.",
        styles['BodyText']
    ))
    story.append(Spacer(1, 12))

    # Table with data
    story.append(Paragraph("Financial Performance", styles['Heading2']))
    data = [
        ['Department', 'Q2 Revenue', 'Q3 Revenue', 'Growth'],
        ['Sales', '$1.2M', '$1.5M', '25%'],
        ['Marketing', '$800K', '$950K', '18.75%'],
        ['Engineering', '$2.1M', '$2.4M', '14.29%'],
        ['Total', '$4.1M', '$4.85M', '18.29%']
    ]

    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 14),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(table)
    story.append(Spacer(1, 12))

    # Bullet list
    story.append(Paragraph("Key Achievements", styles['Heading2']))
    achievements = [
        "Launched new product line with 15% market adoption",
        "Expanded team by 20 employees across all departments",
        "Improved customer satisfaction score to 4.8/5.0",
        "Reduced operational costs by 12% through process optimization"
    ]

    for achievement in achievements:
        story.append(Paragraph(f"• {achievement}", styles['BodyText']))
        story.append(Spacer(1, 6))

    # Build PDF
    doc.build(story)

    # Save metadata
    metadata = {
        "filename": "business_report_01.pdf",
        "type": "pdf",
        "category": "clean_digital",
        "size_bytes": os.path.getsize(output_file),
        "page_count": 1,
        "features": ["tables", "headings", "bullet_lists"],
        "expected_processing_time_seconds": 15,
        "expected_quality_score": 0.95,
        "notes": "Clean digital PDF with table and formatted sections"
    }
    save_metadata(str(output_file), metadata)
    print(f"✓ Generated: {output_file}")


def generate_docx_with_tables():
    """Generate DOCX with complex tables."""
    if not PYTHON_DOCX_AVAILABLE:
        print("Skipping DOCX generation - python-docx not available")
        return

    output_dir = BASE_DIR / "docx" / "with_tables"
    output_file = output_dir / "project_plan_01.docx"

    doc = Document()

    # Title
    title = doc.add_heading('Project Implementation Plan', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Introduction
    doc.add_heading('Project Overview', level=1)
    doc.add_paragraph(
        'This document outlines the implementation plan for the Workshop Document '
        'Processor project. The project aims to deliver a production-ready system '
        'by October 17, 2025.'
    )

    # Table with project timeline
    doc.add_heading('Project Timeline', level=2)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Light Grid Accent 1'

    # Header row
    header_cells = table.rows[0].cells
    header_cells[0].text = 'Phase'
    header_cells[1].text = 'Tasks'
    header_cells[2].text = 'Duration'
    header_cells[3].text = 'Status'

    # Data rows
    data = [
        ('Planning', 'Requirements gathering, Architecture design', '3 days', 'Complete'),
        ('Development', 'Core features implementation', '7 days', 'Complete'),
        ('Testing', 'Integration and load testing', '2 days', 'In Progress'),
        ('Deployment', 'Production deployment and monitoring', '1 day', 'Pending')
    ]

    for phase, tasks, duration, status in data:
        row_cells = table.add_row().cells
        row_cells[0].text = phase
        row_cells[1].text = tasks
        row_cells[2].text = duration
        row_cells[3].text = status

    # Complex table with merged cells
    doc.add_heading('Resource Allocation', level=2)
    doc.add_paragraph(
        'The following table shows resource allocation across team members:'
    )

    resource_table = doc.add_table(rows=5, cols=3)
    resource_table.style = 'Medium Grid 3 Accent 1'

    # Headers
    headers = resource_table.rows[0].cells
    headers[0].text = 'Team Member'
    headers[1].text = 'Role'
    headers[2].text = 'Allocation %'

    # Data
    resources = [
        ('Alice Johnson', 'Product Manager', '50%'),
        ('Bob Smith', 'Lead Developer', '100%'),
        ('Carol Davis', 'QA Engineer', '75%'),
        ('David Lee', 'DevOps Engineer', '60%')
    ]

    for idx, (name, role, allocation) in enumerate(resources, start=1):
        cells = resource_table.rows[idx].cells
        cells[0].text = name
        cells[1].text = role
        cells[2].text = allocation

    # Save document
    doc.save(str(output_file))

    # Save metadata
    metadata = {
        "filename": "project_plan_01.docx",
        "type": "docx",
        "category": "with_tables",
        "size_bytes": os.path.getsize(output_file),
        "page_count": 2,
        "features": ["tables", "headings", "formatted_text", "merged_cells"],
        "expected_processing_time_seconds": 20,
        "expected_quality_score": 0.90,
        "notes": "DOCX with multiple tables including complex formatting"
    }
    save_metadata(str(output_file), metadata)
    print(f"✓ Generated: {output_file}")


def generate_pptx_complex():
    """Generate PPTX with complex layouts."""
    if not PYTHON_PPTX_AVAILABLE:
        print("Skipping PPTX generation - python-pptx not available")
        return

    output_dir = BASE_DIR / "pptx" / "complex_layouts"
    output_file = output_dir / "workshop_presentation_01.pptx"

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Workshop Document Processor"
    subtitle.text = "October 2025 Workshop\nTechnical Overview"

    # Bullet slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Key Features"
    tf = body_shape.text_frame
    tf.text = "Upload and process multiple document formats"

    p = tf.add_paragraph()
    p.text = "PDF, DOCX, PPTX, XLSX support"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "OCR for scanned documents"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Quality and Fast processing modes"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Markdown output optimized for RAG"
    p.level = 0

    # Two-column layout slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    left = top = PptxInches(0.5)
    width = PptxInches(9)
    height = PptxInches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "System Architecture"
    tf.paragraphs[0].font.size = PptxPt(32)
    tf.paragraphs[0].font.bold = True

    # Left column
    left = PptxInches(0.5)
    top = PptxInches(2)
    width = PptxInches(4)
    height = PptxInches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = "Frontend Components:\n\n"
    p = tf.add_paragraph()
    p.text = "• Next.js 14 App Router"
    p = tf.add_paragraph()
    p.text = "• React with TypeScript"
    p = tf.add_paragraph()
    p.text = "• TailwindCSS & shadcn/ui"
    p = tf.add_paragraph()
    p.text = "• Responsive mobile design"

    # Right column
    left = PptxInches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = "Backend Components:\n\n"
    p = tf.add_paragraph()
    p.text = "• FastAPI Python framework"
    p = tf.add_paragraph()
    p.text = "• Docling document processing"
    p = tf.add_paragraph()
    p.text = "• Supabase PostgreSQL & Storage"
    p = tf.add_paragraph()
    p.text = "• DigitalOcean App Platform"

    # Save presentation
    prs.save(str(output_file))

    # Save metadata
    metadata = {
        "filename": "workshop_presentation_01.pptx",
        "type": "pptx",
        "category": "complex_layouts",
        "size_bytes": os.path.getsize(output_file),
        "slide_count": 3,
        "features": ["multi_column", "bullet_lists", "text_boxes", "headings"],
        "expected_processing_time_seconds": 30,
        "expected_quality_score": 0.85,
        "notes": "PPTX with complex multi-column layouts and text boxes"
    }
    save_metadata(str(output_file), metadata)
    print(f"✓ Generated: {output_file}")


def generate_xlsx_with_sheets():
    """Generate XLSX with multiple sheets and formulas."""
    if not OPENPYXL_AVAILABLE:
        print("Skipping XLSX generation - openpyxl not available")
        return

    output_dir = BASE_DIR / "xlsx" / "multiple_sheets"
    output_file = output_dir / "financial_analysis_01.xlsx"

    wb = Workbook()

    # Sales Data sheet
    ws1 = wb.active
    ws1.title = "Sales Data"

    # Headers
    headers = ['Month', 'Product A', 'Product B', 'Product C', 'Total']
    for col, header in enumerate(headers, start=1):
        cell = ws1.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')

    # Data with formulas
    months = ['January', 'February', 'March', 'April', 'May', 'June']
    for row, month in enumerate(months, start=2):
        ws1.cell(row=row, column=1, value=month)
        ws1.cell(row=row, column=2, value=random.randint(1000, 5000))
        ws1.cell(row=row, column=3, value=random.randint(800, 4000))
        ws1.cell(row=row, column=4, value=random.randint(1200, 4500))
        # Formula for total
        ws1.cell(row=row, column=5, value=f"=SUM(B{row}:D{row})")

    # Summary sheet
    ws2 = wb.create_sheet(title="Summary")
    ws2['A1'] = "Summary Statistics"
    ws2['A1'].font = Font(bold=True, size=14)

    ws2['A3'] = "Metric"
    ws2['B3'] = "Value"
    ws2['A3'].font = Font(bold=True)
    ws2['B3'].font = Font(bold=True)

    ws2['A4'] = "Total Sales (Product A)"
    ws2['B4'] = "=SUM('Sales Data'!B2:B7)"

    ws2['A5'] = "Total Sales (Product B)"
    ws2['B5'] = "=SUM('Sales Data'!C2:C7)"

    ws2['A6'] = "Total Sales (Product C)"
    ws2['B6'] = "=SUM('Sales Data'!D2:D7)"

    ws2['A7'] = "Grand Total"
    ws2['B7'] = "=SUM(B4:B6)"
    ws2['A7'].font = Font(bold=True)
    ws2['B7'].font = Font(bold=True)

    # Growth Analysis sheet
    ws3 = wb.create_sheet(title="Growth Analysis")
    ws3['A1'] = "Month-over-Month Growth"
    ws3['A1'].font = Font(bold=True, size=14)

    ws3['A3'] = "Month"
    ws3['B3'] = "Growth %"
    ws3['A3'].font = Font(bold=True)
    ws3['B3'].font = Font(bold=True)

    for row, month in enumerate(months[1:], start=4):
        ws3.cell(row=row, column=1, value=month)
        # Calculate growth percentage
        prev_row = row + 1
        ws3.cell(row=row, column=2, value=f"=(('Sales Data'!E{prev_row}-'Sales Data'!E{prev_row-1})/'Sales Data'!E{prev_row-1})*100")

    # Save workbook
    wb.save(str(output_file))

    # Save metadata
    metadata = {
        "filename": "financial_analysis_01.xlsx",
        "type": "xlsx",
        "category": "multiple_sheets",
        "size_bytes": os.path.getsize(output_file),
        "sheet_count": 3,
        "features": ["multiple_sheets", "formulas", "tables", "cross_sheet_references"],
        "expected_processing_time_seconds": 25,
        "expected_quality_score": 0.88,
        "notes": "XLSX with multiple sheets, formulas, and cross-sheet references"
    }
    save_metadata(str(output_file), metadata)
    print(f"✓ Generated: {output_file}")


def generate_edge_case_files():
    """Generate edge case test files."""

    # File at size limit (9.9MB - under limit)
    size_limits_dir = BASE_DIR / "edge_cases" / "size_limits"
    file_under_limit = size_limits_dir / "file_9.9mb.txt"

    # Generate 9.9MB of text
    chunk = "This is test content for file size validation. " * 100
    target_size = int(9.9 * 1024 * 1024)  # 9.9MB

    with open(file_under_limit, 'w') as f:
        written = 0
        while written < target_size:
            f.write(chunk)
            written += len(chunk.encode('utf-8'))

    metadata = {
        "filename": "file_9.9mb.txt",
        "type": "txt",
        "category": "size_limits",
        "size_bytes": os.path.getsize(file_under_limit),
        "expected_behavior": "should_process_successfully",
        "notes": "File just under 10MB limit - should process successfully"
    }
    save_metadata(str(file_under_limit), metadata)
    print(f"✓ Generated: {file_under_limit}")

    # File over size limit (10.1MB - over limit)
    file_over_limit = size_limits_dir / "file_10.1mb.txt"
    target_size = int(10.1 * 1024 * 1024)  # 10.1MB

    with open(file_over_limit, 'w') as f:
        written = 0
        while written < target_size:
            f.write(chunk)
            written += len(chunk.encode('utf-8'))

    metadata = {
        "filename": "file_10.1mb.txt",
        "type": "txt",
        "category": "size_limits",
        "size_bytes": os.path.getsize(file_over_limit),
        "expected_behavior": "should_fail_with_FILE_TOO_LARGE",
        "expected_error_code": "FILE_TOO_LARGE",
        "notes": "File over 10MB limit - should fail with size error"
    }
    save_metadata(str(file_over_limit), metadata)
    print(f"✓ Generated: {file_over_limit}")

    # File with special characters in filename
    special_chars_dir = BASE_DIR / "edge_cases" / "special_characters"

    # Unicode filename
    unicode_file = special_chars_dir / "测试文档_test_doc_01.txt"
    with open(unicode_file, 'w', encoding='utf-8') as f:
        f.write("Test document with Unicode characters in filename.\n")
        f.write("这是一个测试文档。\n")

    metadata = {
        "filename": "测试文档_test_doc_01.txt",
        "type": "txt",
        "category": "special_characters",
        "size_bytes": os.path.getsize(unicode_file),
        "expected_behavior": "should_process_with_sanitized_filename",
        "notes": "Filename with Chinese Unicode characters"
    }
    save_metadata(str(unicode_file), metadata)
    print(f"✓ Generated: {unicode_file}")

    # Filename with spaces and symbols
    symbols_file = special_chars_dir / "Document (v2.1) - Final [APPROVED].txt"
    with open(symbols_file, 'w') as f:
        f.write("Test document with spaces and symbols in filename.\n")

    metadata = {
        "filename": "Document (v2.1) - Final [APPROVED].txt",
        "type": "txt",
        "category": "special_characters",
        "size_bytes": os.path.getsize(symbols_file),
        "expected_behavior": "should_process_with_sanitized_filename",
        "notes": "Filename with spaces, parentheses, brackets, and symbols"
    }
    save_metadata(str(symbols_file), metadata)
    print(f"✓ Generated: {symbols_file}")

    # Corrupted file (invalid content)
    corrupted_dir = BASE_DIR / "edge_cases" / "corrupted"
    corrupted_pdf = corrupted_dir / "corrupted_document.pdf"

    # Write invalid PDF content
    with open(corrupted_pdf, 'w') as f:
        f.write("This is not a valid PDF file content")

    metadata = {
        "filename": "corrupted_document.pdf",
        "type": "pdf",
        "category": "corrupted",
        "size_bytes": os.path.getsize(corrupted_pdf),
        "expected_behavior": "should_fail_with_CORRUPTED_FILE",
        "expected_error_code": "CORRUPTED_FILE",
        "notes": "Invalid PDF content - should fail gracefully"
    }
    save_metadata(str(corrupted_pdf), metadata)
    print(f"✓ Generated: {corrupted_pdf}")


def generate_multilingual_files():
    """Generate multilingual test files."""
    multilingual_dir = BASE_DIR / "multilingual"

    # Chinese document
    chinese_file = multilingual_dir / "chinese_document.txt"
    with open(chinese_file, 'w', encoding='utf-8') as f:
        f.write("# 中文测试文档\n\n")
        f.write("## 简介\n\n")
        f.write("这是一个用于测试中文文档处理的测试文件。文档处理系统应该能够正确处理和提取中文文本。\n\n")
        f.write("## 测试内容\n\n")
        f.write("1. 中文字符识别\n")
        f.write("2. 标题层次结构\n")
        f.write("3. 文本编码处理\n")

    metadata = {
        "filename": "chinese_document.txt",
        "type": "txt",
        "category": "multilingual",
        "language": "Chinese (Simplified)",
        "size_bytes": os.path.getsize(chinese_file),
        "expected_behavior": "should_process_successfully",
        "notes": "Chinese language document for encoding and text extraction testing"
    }
    save_metadata(str(chinese_file), metadata)
    print(f"✓ Generated: {chinese_file}")

    # Arabic document
    arabic_file = multilingual_dir / "arabic_document.txt"
    with open(arabic_file, 'w', encoding='utf-8') as f:
        f.write("# وثيقة اختبار عربية\n\n")
        f.write("## مقدمة\n\n")
        f.write("هذا ملف اختبار لمعالجة المستندات العربية. يجب أن يكون النظام قادرًا على معالجة واستخراج النص العربي بشكل صحيح.\n\n")
        f.write("## محتوى الاختبار\n\n")
        f.write("1. التعرف على الأحرف العربية\n")
        f.write("2. هيكل العناوين\n")
        f.write("3. معالجة الترميز\n")

    metadata = {
        "filename": "arabic_document.txt",
        "type": "txt",
        "category": "multilingual",
        "language": "Arabic",
        "size_bytes": os.path.getsize(arabic_file),
        "expected_behavior": "should_process_successfully",
        "notes": "Arabic language document for RTL text and encoding testing"
    }
    save_metadata(str(arabic_file), metadata)
    print(f"✓ Generated: {arabic_file}")

    # Spanish document
    spanish_file = multilingual_dir / "spanish_document.txt"
    with open(spanish_file, 'w', encoding='utf-8') as f:
        f.write("# Documento de Prueba en Español\n\n")
        f.write("## Introducción\n\n")
        f.write("Este es un archivo de prueba para el procesamiento de documentos en español. "
                "El sistema debe poder procesar y extraer correctamente el texto en español con caracteres especiales.\n\n")
        f.write("## Contenido de Prueba\n\n")
        f.write("1. Reconocimiento de caracteres españoles (ñ, á, é, í, ó, ú, ü)\n")
        f.write("2. Estructura de encabezados\n")
        f.write("3. Procesamiento de codificación\n")
        f.write("4. Puntuación española (¿interrogación? ¡exclamación!)\n")

    metadata = {
        "filename": "spanish_document.txt",
        "type": "txt",
        "category": "multilingual",
        "language": "Spanish",
        "size_bytes": os.path.getsize(spanish_file),
        "expected_behavior": "should_process_successfully",
        "notes": "Spanish language document with special characters (ñ, accents)"
    }
    save_metadata(str(spanish_file), metadata)
    print(f"✓ Generated: {spanish_file}")


def main():
    """Generate all test documents."""
    print("Generating diverse test documents...\n")

    # Generate documents
    generate_clean_pdf()
    generate_docx_with_tables()
    generate_pptx_complex()
    generate_xlsx_with_sheets()
    generate_edge_case_files()
    generate_multilingual_files()

    print("\n✅ Test document generation complete!")
    print("\nNote: Some documents require additional libraries:")
    print("  - PDF generation: pip install reportlab")
    print("  - DOCX generation: pip install python-docx")
    print("  - PPTX generation: pip install python-pptx")
    print("  - XLSX generation: pip install openpyxl")


if __name__ == "__main__":
    main()
